#!/usr/bin/env python3
from flask import Flask, request, send_file, jsonify
from flask_cors import CORS, cross_origin
from io import BytesIO
import time
import random
import requests
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os
import sys
from typing import List, Optional
from dotenv import load_dotenv

# Load environment variables
load_dotenv()

# Configuration
PRIMARY_MODEL = "gemini-2.0-flash-thinking-exp-01-21"
PEXELS_API_KEY = os.getenv("PEXELS_API_KEY")
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
BASE_RPM = 30
MAX_WORDS = 200
MAX_IMAGES = 3

# Professional Design Settings
DESIGN_CONFIG = {
    "title_font": 'Calibri Light',
    "body_font": 'Calibri',
    "colors": {
        "primary": RGBColor(13, 27, 42),
        "accent": RGBColor(231, 76, 60),
        "neutral": RGBColor(241, 241, 241)
    },
    "slide_size": (Inches(13.33), Inches(7.5)),
    "margins": {
        "side": Inches(0.7),
        "top": Inches(0.3),
        "gap": Inches(0.5)
    },
    "title_size": Pt(36),
    "body_size": Pt(18),
    "max_bullets": 5,
    "text_ratio": 0.55,
    "image_ratio": 0.4,
    "max_lines": 10,
    "chars_per_line": 85
}

PITCH_SECTIONS = [
    "Problem Analysis",
    "Innovative Solution",
    "Market Potential",
    "Business Model",
    "Technology Stack",
    "Go-to-Market Strategy",
    "Financial Projections",
    "Team & Advisors",
    "Competitive Differentiation"
]

class ProfessionalDeck:
    def __init__(self):
        self.prs = Presentation()
        self._setup_master()
    
    def _setup_master(self):
        self.prs.slide_width, self.prs.slide_height = DESIGN_CONFIG["slide_size"]
        master = self.prs.slide_masters[0]
        master.background.fill.solid()
        master.background.fill.fore_color.rgb = DESIGN_CONFIG["colors"]["neutral"]
    
    def add_slide(self, title: str, content: List[str], image_url: Optional[str] = None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        available_width = DESIGN_CONFIG["slide_size"][0] - (2 * DESIGN_CONFIG["margins"]["side"])
        text_width = available_width * DESIGN_CONFIG["text_ratio"]
        image_width = available_width * DESIGN_CONFIG["image_ratio"] if image_url else 0
        
        self._add_title(slide, title, available_width)
        self._add_content(slide, content, text_width)
        if image_url:
            self._add_image(slide, image_url, text_width, image_width)

    def _add_title(self, slide, title: str, width: Inches):
        title_box = slide.shapes.add_textbox(
            left=DESIGN_CONFIG["margins"]["side"],
            top=DESIGN_CONFIG["margins"]["top"],
            width=width,
            height=Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = DESIGN_CONFIG["title_font"]
        p.font.size = DESIGN_CONFIG["title_size"]
        p.font.color.rgb = DESIGN_CONFIG["colors"]["accent"]
        p.alignment = PP_ALIGN.LEFT

    def _add_content(self, slide, content: List[str], width: Inches):
        textbox = slide.shapes.add_textbox(
            left=DESIGN_CONFIG["margins"]["side"],
            top=DESIGN_CONFIG["margins"]["top"] + Inches(0.9),
            width=width,
            height=Inches(4.2)
        )
        tf = textbox.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        
        current_lines = 0
        for idx, point in enumerate(content[:DESIGN_CONFIG["max_bullets"]]):
            if current_lines >= DESIGN_CONFIG["max_lines"]:
                break
                
            truncated = self._truncate_text(
                point, 
                DESIGN_CONFIG["chars_per_line"],
                DESIGN_CONFIG["max_lines"] - current_lines
            )
            
            p = tf.add_paragraph()
            p.text = truncated
            p.font.name = DESIGN_CONFIG["body_font"]
            p.font.size = DESIGN_CONFIG["body_size"]
            p.space_after = Pt(10)
            p.level = 0 if point.startswith("•") else 1
            p.font.color.rgb = DESIGN_CONFIG["colors"]["primary"]
            
            line_count = len(truncated) // DESIGN_CONFIG["chars_per_line"] + 1
            current_lines += line_count
            
            if idx > 0:
                p.space_before = Pt(8)

    def _truncate_text(self, text: str, chars_per_line: int, max_lines: int) -> str:
        max_chars = chars_per_line * max_lines
        truncated = text[:max_chars]
        if len(text) > max_chars:
            return truncated.rsplit(' ', 1)[0] + '...'
        return truncated

    def _add_image(self, slide, image_url: str, text_width: Inches, image_width: Inches):
        try:
            response = requests.get(image_url, timeout=10)
            img_stream = BytesIO(response.content)
            image_left = DESIGN_CONFIG["margins"]["side"] + text_width + DESIGN_CONFIG["margins"]["gap"]
            
            slide.shapes.add_picture(
                img_stream,
                left=image_left,
                top=DESIGN_CONFIG["margins"]["top"] + Inches(0.9),
                width=image_width
            )
        except Exception as e:
            print(f"Image error: {str(e)}")

class ContentGenerator:
    def __init__(self):
        genai.configure(api_key=GEMINI_API_KEY)
        self.model = genai.GenerativeModel(PRIMARY_MODEL)
        self.last_request = 0
        self.main_theme = ""
        self.used_images = set()

    def _rate_limit(self):
        elapsed = time.time() - self.last_request
        if elapsed < 60 / BASE_RPM:
            time.sleep(60 / BASE_RPM - elapsed)

    def _get_thematic_image(self, section: str) -> Optional[str]:
        try:
            query = f"{self.main_theme} {section} professional business"
            response = requests.get(
                "https://api.pexels.com/v1/search",
                params={
                    "query": query,
                    "per_page": MAX_IMAGES,
                    "orientation": "landscape"
                },
                headers={"Authorization": PEXELS_API_KEY},
                timeout=10
            )
            photos = response.json().get("photos", [])
            if photos:
                for photo in random.sample(photos, min(len(photos), MAX_IMAGES)):
                    img_url = photo["src"]["large"]
                    if img_url not in self.used_images:
                        self.used_images.add(img_url)
                        return img_url
                return photos[0]["src"]["large"]
        except Exception as e:
            print(f"Image search failed: {str(e)}")
        return None

    def generate_slide_content(self, section: str) -> tuple:
        self._rate_limit()
        prompt = f"""Create {section} section for {self.main_theme} with:
        - MAX {MAX_WORDS} words
        - 3-5 concise bullet points
        - Industry-specific metrics
        - Strictly no markdown"""
        
        try:
            response = self.model.generate_content(prompt)
            content = self._process_text(response.text)
            image_url = self._get_thematic_image(section)
            self.last_request = time.time()
            return content, image_url
        except Exception as e:
            return [f"Content error: {str(e)}"], None

    def _process_text(self, text: str) -> List[str]:
        clean_lines = []
        for line in text.split('\n'):
            line = line.strip()
            if line and not any(term in line for term in ["TBD", "N/A"]):
                prefix = "• " if not line.startswith("•") else ""
                clean_lines.append(f"{prefix}{line}")
        return clean_lines[:DESIGN_CONFIG["max_bullets"]]

# Flask Server Setup
app = Flask(__name__)
CORS(app)

def validate_env():
    missing = [k for k in ["GEMINI_API_KEY", "PEXELS_API_KEY"] if not os.getenv(k)]
    if missing:
        raise EnvironmentError(f"Missing API keys: {', '.join(missing)}")

@app.route('/generate', methods=['POST'])
@cross_origin()
def generate_deck():
    try:
        validate_env()
        data = request.get_json()
        
        if not data or 'idea' not in data:
            return jsonify({"error": "Missing 'idea' in request"}), 400
            
        idea = data['idea'].strip()
        if not idea:
            return jsonify({"error": "Empty idea parameter"}), 400

        deck = ProfessionalDeck()
        generator = ContentGenerator()
        generator.main_theme = idea
        
        # Title slide
        deck.add_slide(idea, [
            f"Strategic Presentation: {idea}",
            f"Generated: {time.strftime('%Y-%m-%d')}",
            "Confidential & Proprietary"
        ])
        
        # Content slides
        for section in PITCH_SECTIONS:
            content, image_url = generator.generate_slide_content(section)
            deck.add_slide(section, content, image_url)
        
        # Save to in-memory buffer
        buffer = BytesIO()
        deck.prs.save(buffer)
        buffer.seek(0)
        
        return send_file(
            buffer,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name=f"PitchDeck_{idea[:50].replace(' ', '_')}.pptx"
        )
        
    except EnvironmentError as e:
        return jsonify({"error": str(e)}), 500
    except Exception as e:
        return jsonify({"error": f"Generation failed: {str(e)}"}), 500

if __name__ == '__main__':
    port = int(os.getenv("PORT", 5000))
    app.run(debug=False, host="0.0.0.0", port=port)
